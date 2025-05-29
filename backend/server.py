from fastapi import FastAPI, APIRouter, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any
import uuid
from datetime import datetime, timedelta
import yfinance as yf
import pandas as pd
import numpy as np
from openpyxl import Workbook
from openpyxl.chart import LineChart, Reference
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import seaborn as sns
import io
import base64
from scipy import stats
import json

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Create the main app without a prefix
app = FastAPI(title="FCN Investment Analysis Calculator")

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")

# Ensure reports directory exists
REPORTS_DIR = ROOT_DIR / "reports"
REPORTS_DIR.mkdir(exist_ok=True)

# Define Models
class StockInfo(BaseModel):
    symbol: str
    name: str
    current_price: float
    market_cap: Optional[float] = None
    pe_ratio: Optional[float] = None
    dividend_yield: Optional[float] = None
    beta: Optional[float] = None
    exchange: str

class FCNParameters(BaseModel):
    coupon_rate: float  # Annual coupon rate (e.g., 5.5 for 5.5%)
    face_value: float  # Face value of the note
    maturity_months: int  # Maturity in months (FCNs typically use monthly terms)
    reference_prices: Dict[str, float]  # Reference prices for each stock in the basket
    strike_prices: Dict[str, float]  # Strike prices for each stock (often same as reference prices)
    put_strike_prices: Dict[str, float]  # Put strike prices for each stock (for equity conversion)
    knock_out_barrier_pct: float  # Knock-out barrier as % of reference price (e.g., 110.0 for 110%)
    knock_in_barrier_pct: float  # Knock-in barrier as % of reference price (e.g., 70.0 for 70%)
    barrier_style: str = "american"  # "american" (continuous monitoring) or "european" (observation dates only)
    observation_frequency: str = "monthly"  # monthly, weekly, daily
    autocallable: bool = True  # Whether the note can be called early on knock-out
    basket_type: str = "worst_of"  # "worst_of" (standard FCN), "best_of", "average"

class FCNAnalysisRequest(BaseModel):
    symbols: List[str]  # List of stock symbols
    fcn_params: FCNParameters
    analysis_period: int = 252  # Trading days for historical analysis
    scenarios: Dict[str, float] = Field(default_factory=lambda: {
        "base_case": 0.0,
        "bull_case": 0.15,
        "bear_case": -0.20
    })

class FCNAnalysisResult(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    request_params: FCNAnalysisRequest
    stocks_info: List[StockInfo]
    fcn_metrics: Dict[str, Any]
    scenario_analysis: Dict[str, Any]
    risk_metrics: Dict[str, Any]
    charts: Dict[str, str]  # Base64 encoded charts
    created_at: datetime = Field(default_factory=datetime.utcnow)

class ReportRequest(BaseModel):
    analysis_id: str
    report_type: str  # "excel" or "powerpoint"
    include_charts: bool = True

# FCN Calculation Functions
def calculate_basket_performance(current_prices: Dict[str, float], reference_prices: Dict[str, float], 
                                basket_type: str = "worst_of") -> Dict[str, float]:
    """Calculate basket performance based on basket type (worst-of, best-of, average)"""
    performances = {}
    for symbol in current_prices.keys():
        if symbol in reference_prices:
            performances[symbol] = (current_prices[symbol] / reference_prices[symbol] - 1) * 100
    
    if basket_type == "worst_of":
        worst_symbol = min(performances, key=performances.get)
        return {
            "performance": performances[worst_symbol],
            "worst_symbol": worst_symbol,
            "worst_performance": performances[worst_symbol],
            "all_performances": performances
        }
    elif basket_type == "best_of":
        best_symbol = max(performances, key=performances.get)
        return {
            "performance": performances[best_symbol],
            "best_symbol": best_symbol,
            "best_performance": performances[best_symbol],
            "all_performances": performances
        }
    else:  # average
        avg_performance = sum(performances.values()) / len(performances)
        return {
            "performance": avg_performance,
            "average_performance": avg_performance,
            "all_performances": performances
        }

def check_basket_barriers(current_prices: Dict[str, float], reference_prices: Dict[str, float],
                         knock_out_barrier_pct: float, knock_in_barrier_pct: float) -> Dict[str, bool]:
    """Check if any stock in the basket has breached barriers"""
    knocked_out = False
    knocked_in = False
    barrier_events = {}
    
    for symbol, current_price in current_prices.items():
        if symbol in reference_prices:
            ref_price = reference_prices[symbol]
            ko_barrier = ref_price * (knock_out_barrier_pct / 100)
            ki_barrier = ref_price * (knock_in_barrier_pct / 100)
            
            stock_ko = current_price >= ko_barrier
            stock_ki = current_price <= ki_barrier
            
            barrier_events[symbol] = {
                "knocked_out": stock_ko,
                "knocked_in": stock_ki,
                "ko_barrier": ko_barrier,
                "ki_barrier": ki_barrier
            }
            
            # If ANY stock breaches barriers, the entire FCN is affected
            if stock_ko:
                knocked_out = True
            if stock_ki:
                knocked_in = True
    
    return {
        "knocked_out": knocked_out,
        "knocked_in": knocked_in,
        "barrier_events": barrier_events
    }

def calculate_fcn_payoff(final_prices: Dict[str, float], reference_prices: Dict[str, float], 
                        strike_prices: Dict[str, float], put_strike_prices: Dict[str, float],
                        knock_out_barrier_pct: float, knock_in_barrier_pct: float, 
                        coupon_rate: float, face_value: float, maturity_months: int, 
                        barrier_breached: Dict[str, bool], basket_type: str = "worst_of", 
                        early_redemption_month: Optional[int] = None) -> Dict[str, float]:
    """Calculate FCN payoff based on basket structure with worst-of performance and put strikes"""
    
    # Monthly coupon
    monthly_coupon = face_value * (coupon_rate / 100) / 12
    
    if early_redemption_month:
        # Early redemption due to knock-out
        total_coupons = monthly_coupon * early_redemption_month
        return {
            "payoff": face_value + total_coupons,
            "total_return": (total_coupons) / face_value * 100,
            "coupons_received": total_coupons,
            "redemption_type": "early_knockout"
        }
    
    # Calculate total coupons for full term
    total_coupons = monthly_coupon * maturity_months
    
    if not barrier_breached["knock_in"]:
        # No knock-in occurred, receive face value + all coupons
        return {
            "payoff": face_value + total_coupons,
            "total_return": (total_coupons) / face_value * 100,
            "coupons_received": total_coupons,
            "redemption_type": "full_term_protected"
        }
    else:
        # Knock-in occurred, equity exposure based on worst-performing stock vs its PUT STRIKE
        basket_perf = calculate_basket_performance(final_prices, reference_prices, basket_type)
        
        if basket_type == "worst_of":
            worst_symbol = basket_perf["worst_symbol"]
            put_strike = put_strike_prices[worst_symbol]
            worst_final = final_prices[worst_symbol]
            
            # Use put strike for equity exposure calculation
            equity_performance = worst_final / put_strike
        else:
            # For other basket types, use average performance vs average put strike
            avg_performance = basket_perf["performance"] / 100 + 1  # Convert % to multiplier
            equity_performance = avg_performance
        
        equity_payoff = face_value * equity_performance
        
        return {
            "payoff": equity_payoff + total_coupons,
            "total_return": (equity_payoff + total_coupons - face_value) / face_value * 100,
            "coupons_received": total_coupons,
            "equity_performance": (equity_performance - 1) * 100,
            "worst_performer": basket_perf.get("worst_symbol", "N/A"),
            "put_strike_used": put_strike_prices.get(basket_perf.get("worst_symbol", ""), 0),
            "redemption_type": "equity_exposure"
        }

def calculate_fcn_metrics(current_prices: Dict[str, float], fcn_params: FCNParameters) -> Dict[str, Dict[str, float]]:
    """Calculate FCN-specific metrics for basket structure"""
    
    metrics = {}
    
    # Calculate basket performance
    basket_perf = calculate_basket_performance(current_prices, fcn_params.reference_prices, fcn_params.basket_type)
    
    # Check barriers for all stocks
    barrier_status = check_basket_barriers(
        current_prices, fcn_params.reference_prices,
        fcn_params.knock_out_barrier_pct, fcn_params.knock_in_barrier_pct
    )
    
    for symbol in current_prices.keys():
        if symbol in fcn_params.reference_prices:
            current_price = current_prices[symbol]
            ref_price = fcn_params.reference_prices[symbol]
            
            # Calculate actual barrier levels
            knock_out_barrier = ref_price * (fcn_params.knock_out_barrier_pct / 100)
            knock_in_barrier = ref_price * (fcn_params.knock_in_barrier_pct / 100)
            
            # Distance to barriers from current price
            distance_to_knockout = ((knock_out_barrier - current_price) / current_price) * 100
            distance_to_knockin = ((current_price - knock_in_barrier) / current_price) * 100
            
            # Performance vs reference price
            performance_vs_reference = (current_price / ref_price - 1) * 100
            
            # Moneyness relative to strike and put strike
            strike_price = fcn_params.strike_prices.get(symbol, ref_price)
            put_strike_price = fcn_params.put_strike_prices.get(symbol, ref_price)
            moneyness = (current_price / strike_price - 1) * 100
            put_moneyness = (current_price / put_strike_price - 1) * 100
            
            metrics[symbol] = {
                "current_yield": fcn_params.coupon_rate,
                "distance_to_knockout": distance_to_knockout,
                "distance_to_knockin": distance_to_knockin,
                "performance_vs_reference": performance_vs_reference,
                "moneyness": moneyness,
                "reference_price": ref_price,
                "knockout_barrier": knock_out_barrier,
                "knockin_barrier": knock_in_barrier,
                "knockout_barrier_pct": fcn_params.knock_out_barrier_pct,
                "knockin_barrier_pct": fcn_params.knock_in_barrier_pct,
                "monthly_coupon": fcn_params.face_value * (fcn_params.coupon_rate / 100) / 12 / len(current_prices),  # Split among basket
                "is_worst_performer": symbol == basket_perf.get("worst_symbol", ""),
                "basket_performance": basket_perf["performance"]
            }
    
    return metrics

def monte_carlo_simulation(stock_prices: pd.DataFrame, fcn_params: FCNParameters, 
                          num_simulations: int = 10000) -> Dict[str, Any]:
    """Run Monte Carlo simulation for FCN basket analysis with worst-of performance"""
    
    symbols = list(fcn_params.reference_prices.keys())
    
    # Calculate returns for each stock
    returns_data = {}
    for symbol in symbols:
        if symbol in stock_prices.columns:
            returns_data[symbol] = {
                'returns': stock_prices[symbol].pct_change().dropna(),
                'mu': stock_prices[symbol].pct_change().dropna().mean() * 252,
                'sigma': stock_prices[symbol].pct_change().dropna().std() * np.sqrt(252)
            }
    
    # Monte Carlo simulation
    payoffs = []
    knock_in_events = 0
    knock_out_events = 0
    early_redemptions = []
    worst_performers = []
    
    for sim in range(num_simulations):
        # Generate price paths for all stocks in basket
        days = int(fcn_params.maturity_months * 21)  # Approximate trading days per month
        dt = 1/252  # Daily time step
        
        price_paths = {}
        for symbol in symbols:
            if symbol in returns_data:
                price_paths[symbol] = [fcn_params.reference_prices[symbol]]
            
        knocked_in = False
        knocked_out = False
        redemption_month = None
        
        for day in range(days):
            current_prices = {}
            
            # Generate new prices for each stock
            for symbol in symbols:
                if symbol in returns_data and symbol in price_paths:
                    mu = returns_data[symbol]['mu']
                    sigma = returns_data[symbol]['sigma']
                    random_shock = np.random.normal(0, 1)
                    price_change = mu * dt + sigma * np.sqrt(dt) * random_shock
                    new_price = price_paths[symbol][-1] * np.exp(price_change)
                    price_paths[symbol].append(new_price)
                    current_prices[symbol] = new_price
            
            # Check barriers for the entire basket
            current_month = int(day / 21) + 1  # Current month
            
            if fcn_params.barrier_style == "american" or (
                fcn_params.barrier_style == "european" and day % 21 == 0
            ):
                # Check barriers across all stocks in basket
                barrier_status = check_basket_barriers(
                    current_prices, fcn_params.reference_prices,
                    fcn_params.knock_out_barrier_pct, fcn_params.knock_in_barrier_pct
                )
                
                # Check knock-out barrier
                if barrier_status["knocked_out"] and fcn_params.autocallable:
                    knocked_out = True
                    redemption_month = current_month
                    knock_out_events += 1
                    break
                
                # Check knock-in barrier
                if barrier_status["knocked_in"]:
                    knocked_in = True
                    knock_in_events += 1
        
        # Get final prices
        final_prices = {}
        for symbol in symbols:
            if symbol in price_paths:
                final_prices[symbol] = price_paths[symbol][-1]
        
        # Calculate FCN payoff using basket structure
        payoff_result = calculate_fcn_payoff(
            final_prices=final_prices,
            reference_prices=fcn_params.reference_prices,
            strike_prices=fcn_params.strike_prices,
            knock_out_barrier_pct=fcn_params.knock_out_barrier_pct,
            knock_in_barrier_pct=fcn_params.knock_in_barrier_pct,
            coupon_rate=fcn_params.coupon_rate,
            face_value=fcn_params.face_value,
            maturity_months=fcn_params.maturity_months,
            barrier_breached={"knock_in": knocked_in, "knock_out": knocked_out},
            basket_type=fcn_params.basket_type,
            early_redemption_month=redemption_month
        )
        
        payoffs.append(payoff_result["payoff"])
        worst_performers.append(payoff_result.get("worst_performer", "N/A"))
        
        if redemption_month:
            early_redemptions.append(redemption_month)
    
    # Calculate statistics
    avg_redemption_month = np.mean(early_redemptions) if early_redemptions else fcn_params.maturity_months
    
    # Count worst performer frequency
    worst_performer_counts = {}
    for performer in worst_performers:
        if performer != "N/A":
            worst_performer_counts[performer] = worst_performer_counts.get(performer, 0) + 1
    
    most_frequent_worst = max(worst_performer_counts, key=worst_performer_counts.get) if worst_performer_counts else "N/A"
    
    return {
        'basket_type': fcn_params.basket_type,
        'symbols': symbols,
        'expected_payoff': np.mean(payoffs),
        'payoff_std': np.std(payoffs),
        'knock_in_probability': knock_in_events / num_simulations * 100,
        'knock_out_probability': knock_out_events / num_simulations * 100,
        'avg_redemption_month': avg_redemption_month,
        'var_95': np.percentile(payoffs, 5),
        'var_99': np.percentile(payoffs, 1),
        'max_payoff': np.max(payoffs),
        'min_payoff': np.min(payoffs),
        'worst_performer_frequency': worst_performer_counts,
        'most_frequent_worst': most_frequent_worst
    }

def create_price_chart(stock_prices: pd.DataFrame) -> str:
    """Create price chart and return as base64 string"""
    plt.figure(figsize=(12, 6))
    for symbol in stock_prices.columns:
        plt.plot(stock_prices.index, stock_prices[symbol], label=symbol, linewidth=2)
    
    plt.title('Stock Price History', fontsize=16, fontweight='bold')
    plt.xlabel('Date', fontsize=12)
    plt.ylabel('Price ($)', fontsize=12)
    plt.legend()
    plt.grid(True, alpha=0.3)
    plt.tight_layout()
    
    # Convert to base64
    buffer = io.BytesIO()
    plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
    buffer.seek(0)
    image_base64 = base64.b64encode(buffer.getvalue()).decode()
    plt.close()
    
    return image_base64

def create_payoff_distribution_chart(monte_carlo_results: Dict) -> str:
    """Create payoff distribution chart for basket FCN"""
    fig, ax = plt.subplots(1, 1, figsize=(10, 6))
    
    # Create sample distribution for visualization (simplified)
    expected_payoff = monte_carlo_results['expected_payoff']
    payoff_std = monte_carlo_results['payoff_std']
    sample_payoffs = np.random.normal(expected_payoff, payoff_std, 1000)
    
    ax.hist(sample_payoffs, bins=50, alpha=0.7, edgecolor='black')
    ax.axvline(expected_payoff, color='red', linestyle='--', linewidth=2, 
               label=f'Expected: ${expected_payoff:.2f}')
    ax.axvline(monte_carlo_results['var_95'], color='orange', linestyle='--', 
               label=f'VaR 95%: ${monte_carlo_results["var_95"]:.2f}')
    
    ax.set_title(f'FCN Payoff Distribution - {monte_carlo_results["basket_type"].replace("_", " ").title()} Basket', 
                fontsize=14, fontweight='bold')
    ax.set_xlabel('Payoff ($)', fontsize=12)
    ax.set_ylabel('Frequency', fontsize=12)
    ax.legend()
    ax.grid(True, alpha=0.3)
    
    # Add basket info
    if monte_carlo_results.get('most_frequent_worst') != "N/A":
        ax.text(0.02, 0.98, f'Most Frequent Worst Performer: {monte_carlo_results["most_frequent_worst"]}', 
                transform=ax.transAxes, verticalalignment='top', 
                bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))
    
    plt.tight_layout()
    
    buffer = io.BytesIO()
    plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
    buffer.seek(0)
    image_base64 = base64.b64encode(buffer.getvalue()).decode()
    plt.close()
    
    return image_base64

# API Routes
@api_router.get("/")
async def root():
    return {"message": "FCN Investment Analysis Calculator API"}

@api_router.get("/stock/{symbol}")
async def get_stock_info(symbol: str):
    """Get current stock information"""
    try:
        ticker = yf.Ticker(symbol)
        info = ticker.info
        hist = ticker.history(period="1d")
        
        if hist.empty:
            raise HTTPException(status_code=404, detail=f"Stock symbol {symbol} not found")
        
        current_price = hist['Close'].iloc[-1]
        
        stock_info = StockInfo(
            symbol=symbol.upper(),
            name=info.get('longName', symbol),
            current_price=float(current_price),
            market_cap=info.get('marketCap'),
            pe_ratio=info.get('forwardPE'),
            dividend_yield=info.get('dividendYield'),
            beta=info.get('beta'),
            exchange=info.get('exchange', 'Unknown')
        )
        
        return stock_info
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error fetching stock data: {str(e)}")

@api_router.post("/analyze", response_model=FCNAnalysisResult)
async def analyze_fcn(request: FCNAnalysisRequest):
    """Perform comprehensive FCN analysis"""
    try:
        # Validate basket structure - FCNs typically require exactly 2 stocks
        if len(request.symbols) < 2:
            raise HTTPException(status_code=400, detail="FCN requires at least 2 underlying stocks for basket structure")
        
        # Fetch stock data
        stocks_info = []
        stock_prices = pd.DataFrame()
        current_prices = {}
        
        for symbol in request.symbols:
            ticker = yf.Ticker(symbol)
            info = ticker.info
            
            # Get historical data
            hist = ticker.history(period=f"{request.analysis_period}d")
            if hist.empty:
                raise HTTPException(status_code=404, detail=f"No data found for {symbol}")
            
            stock_prices[symbol] = hist['Close']
            
            # Get current stock info
            current_price = hist['Close'].iloc[-1]
            current_prices[symbol] = current_price
            
            stock_info = StockInfo(
                symbol=symbol.upper(),
                name=info.get('longName', symbol),
                current_price=float(current_price),
                market_cap=info.get('marketCap'),
                pe_ratio=info.get('forwardPE'),
                dividend_yield=info.get('dividendYield'),
                beta=info.get('beta'),
                exchange=info.get('exchange', 'Unknown')
            )
            stocks_info.append(stock_info)
        
        # Calculate FCN metrics for basket
        fcn_metrics = calculate_fcn_metrics(current_prices, request.fcn_params)
        
        # Monte Carlo simulation with basket structure
        monte_carlo_results = monte_carlo_simulation(stock_prices, request.fcn_params)
        
        # Scenario analysis with basket structure
        scenario_analysis = {}
        for scenario_name, return_pct in request.scenarios.items():
            scenario_results = {}
            
            # Apply scenario return to all stocks in basket based on reference prices
            future_prices = {}
            for symbol in request.symbols:
                ref_price = request.fcn_params.reference_prices.get(symbol, current_prices[symbol])
                future_prices[symbol] = ref_price * (1 + return_pct)
            
            # Calculate barriers for scenario
            barrier_status = check_basket_barriers(
                future_prices, request.fcn_params.reference_prices,
                request.fcn_params.knock_out_barrier_pct, request.fcn_params.knock_in_barrier_pct
            )
            
            # Calculate payoff for scenario
            payoff_result = calculate_fcn_payoff(
                final_prices=future_prices,
                reference_prices=request.fcn_params.reference_prices,
                strike_prices=request.fcn_params.strike_prices,
                knock_out_barrier_pct=request.fcn_params.knock_out_barrier_pct,
                knock_in_barrier_pct=request.fcn_params.knock_in_barrier_pct,
                coupon_rate=request.fcn_params.coupon_rate,
                face_value=request.fcn_params.face_value,
                maturity_months=request.fcn_params.maturity_months,
                barrier_breached={
                    "knock_in": barrier_status["knocked_in"],
                    "knock_out": barrier_status["knocked_out"]
                },
                basket_type=request.fcn_params.basket_type,
                early_redemption_month=None
            )
            
            # Calculate basket performance for scenario
            basket_perf = calculate_basket_performance(
                future_prices, request.fcn_params.reference_prices, request.fcn_params.basket_type
            )
            
            scenario_results = {
                "basket_performance": basket_perf["performance"],
                "worst_performer": basket_perf.get("worst_symbol", "N/A"),
                "payoff": payoff_result["payoff"],
                "total_return": payoff_result["total_return"],
                "coupons_received": payoff_result["coupons_received"],
                "redemption_type": payoff_result["redemption_type"],
                "individual_performances": basket_perf["all_performances"]
            }
            
            scenario_analysis[scenario_name] = scenario_results
        
        # Risk metrics for basket FCN
        risk_metrics = {
            "basket_metrics": {
                "basket_type": monte_carlo_results['basket_type'],
                "expected_payoff": monte_carlo_results['expected_payoff'],
                "payoff_std": monte_carlo_results['payoff_std'],
                "knock_in_probability": monte_carlo_results['knock_in_probability'],
                "knock_out_probability": monte_carlo_results['knock_out_probability'],
                "avg_redemption_month": monte_carlo_results['avg_redemption_month'],
                "var_95": monte_carlo_results['var_95'],
                "var_99": monte_carlo_results['var_99'],
                "max_payoff": monte_carlo_results['max_payoff'],
                "min_payoff": monte_carlo_results['min_payoff'],
                "worst_performer_frequency": monte_carlo_results['worst_performer_frequency'],
                "most_frequent_worst": monte_carlo_results['most_frequent_worst']
            }
        }
        
        # Add individual stock risk metrics
        for symbol in request.symbols:
            if symbol in stock_prices.columns:
                returns = stock_prices[symbol].pct_change().dropna()
                
                risk_metrics[symbol] = {
                    "volatility_annualized": returns.std() * np.sqrt(252) * 100,
                    "sharpe_ratio": (returns.mean() * 252) / (returns.std() * np.sqrt(252)),
                    "max_drawdown": ((stock_prices[symbol] / stock_prices[symbol].cummax()) - 1).min() * 100,
                    "is_worst_performer": fcn_metrics[symbol].get("is_worst_performer", False),
                    "performance_vs_reference": fcn_metrics[symbol].get("performance_vs_reference", 0)
                }
        
        # Create charts
        price_chart = create_price_chart(stock_prices)
        payoff_chart = create_payoff_distribution_chart(monte_carlo_results)
        
        charts = {
            "price_history": price_chart,
            "payoff_distribution": payoff_chart
        }
        
        # Create analysis result
        analysis = FCNAnalysisResult(
            request_params=request,
            stocks_info=stocks_info,
            fcn_metrics=fcn_metrics,
            scenario_analysis=scenario_analysis,
            risk_metrics=risk_metrics,
            charts=charts
        )
        
        # Save to database
        await db.fcn_analyses.insert_one(analysis.dict())
        
        return analysis
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Analysis error: {str(e)}")

@api_router.get("/analysis/{analysis_id}")
async def get_analysis(analysis_id: str):
    """Get saved analysis by ID"""
    analysis = await db.fcn_analyses.find_one({"id": analysis_id})
    if not analysis:
        raise HTTPException(status_code=404, detail="Analysis not found")
    
    return FCNAnalysisResult(**analysis)

@api_router.get("/analyses")
async def list_analyses(limit: int = 20):
    """List recent analyses"""
    analyses = await db.fcn_analyses.find().sort("created_at", -1).limit(limit).to_list(limit)
    return [{"id": a["id"], "created_at": a["created_at"], "symbols": a["request_params"]["symbols"]} for a in analyses]

def create_excel_report(analysis: FCNAnalysisResult, file_path: Path):
    """Create Excel report for FCN basket analysis"""
    wb = Workbook()
    
    # Summary sheet
    ws_summary = wb.active
    ws_summary.title = "Executive Summary"
    
    # Header styling
    header_font = Font(bold=True, size=14)
    header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
    
    # Title
    ws_summary['A1'] = "FCN Basket Investment Analysis Report"
    ws_summary['A1'].font = Font(bold=True, size=16)
    ws_summary.merge_cells('A1:F1')
    
    # Analysis parameters
    row = 3
    ws_summary[f'A{row}'] = "FCN Basket Parameters"
    ws_summary[f'A{row}'].font = header_font
    row += 1
    
    params = analysis.request_params.fcn_params
    ws_summary[f'A{row}'] = "Coupon Rate:"
    ws_summary[f'B{row}'] = f"{params.coupon_rate}%"
    row += 1
    ws_summary[f'A{row}'] = "Face Value:"
    ws_summary[f'B{row}'] = f"${params.face_value:,.2f}"
    row += 1
    ws_summary[f'A{row}'] = "Maturity:"
    ws_summary[f'B{row}'] = f"{params.maturity_months} months"
    row += 1
    ws_summary[f'A{row}'] = "Basket Type:"
    ws_summary[f'B{row}'] = f"{params.basket_type.replace('_', ' ').title()}"
    row += 1
    ws_summary[f'A{row}'] = "Knock-Out Barrier:"
    ws_summary[f'B{row}'] = f"{params.knock_out_barrier_pct}%"
    row += 1
    ws_summary[f'A{row}'] = "Knock-In Barrier:"
    ws_summary[f'B{row}'] = f"{params.knock_in_barrier_pct}%"
    row += 2
    
    # Reference prices
    ws_summary[f'A{row}'] = "Reference Prices"
    ws_summary[f'A{row}'].font = header_font
    row += 1
    
    for symbol, price in params.reference_prices.items():
        ws_summary[f'A{row}'] = f"{symbol}:"
        ws_summary[f'B{row}'] = f"${price:,.2f}"
        row += 1
    row += 1
    
    # Stock information
    ws_summary[f'A{row}'] = "Underlying Stocks"
    ws_summary[f'A{row}'].font = header_font
    row += 1
    
    headers = ["Symbol", "Name", "Current Price", "Market Cap", "P/E Ratio", "Exchange"]
    for col, header in enumerate(headers, 1):
        ws_summary.cell(row=row, column=col, value=header).font = Font(bold=True)
    row += 1
    
    for stock in analysis.stocks_info:
        ws_summary.cell(row=row, column=1, value=stock.symbol)
        ws_summary.cell(row=row, column=2, value=stock.name)
        ws_summary.cell(row=row, column=3, value=f"${stock.current_price:.2f}")
        ws_summary.cell(row=row, column=4, value=stock.market_cap if stock.market_cap else "N/A")
        ws_summary.cell(row=row, column=5, value=stock.pe_ratio if stock.pe_ratio else "N/A")
        ws_summary.cell(row=row, column=6, value=stock.exchange)
        row += 1
    
    # Basket Performance sheet
    ws_basket = wb.create_sheet("Basket Performance")
    row = 1
    ws_basket[f'A{row}'] = "FCN Basket Performance Metrics"
    ws_basket[f'A{row}'].font = Font(bold=True, size=14)
    row += 2
    
    # Basket summary metrics
    basket_metrics = analysis.risk_metrics.get('basket_metrics', {})
    ws_basket[f'A{row}'] = "Expected Payoff:"
    ws_basket[f'B{row}'] = f"${basket_metrics.get('expected_payoff', 0):.2f}"
    row += 1
    ws_basket[f'A{row}'] = "Knock-Out Probability:"
    ws_basket[f'B{row}'] = f"{basket_metrics.get('knock_out_probability', 0):.2f}%"
    row += 1
    ws_basket[f'A{row}'] = "Knock-In Probability:"
    ws_basket[f'B{row}'] = f"{basket_metrics.get('knock_in_probability', 0):.2f}%"
    row += 1
    ws_basket[f'A{row}'] = "Most Frequent Worst Performer:"
    ws_basket[f'B{row}'] = basket_metrics.get('most_frequent_worst', 'N/A')
    row += 2
    
    # Individual stock metrics
    headers = ["Symbol", "Current Price", "Reference Price", "Performance vs Ref", "Worst Performer", "Volatility"]
    for col, header in enumerate(headers, 1):
        ws_basket.cell(row=row, column=col, value=header).font = Font(bold=True)
    row += 1
    
    for symbol, metrics in analysis.fcn_metrics.items():
        ws_basket.cell(row=row, column=1, value=symbol)
        
        # Find stock info for current price
        stock = next((s for s in analysis.stocks_info if s.symbol == symbol), None)
        if stock:
            ws_basket.cell(row=row, column=2, value=f"${stock.current_price:.2f}")
        
        ref_price = params.reference_prices.get(symbol, 0)
        ws_basket.cell(row=row, column=3, value=f"${ref_price:.2f}")
        ws_basket.cell(row=row, column=4, value=f"{metrics.get('performance_vs_reference', 0):.2f}%")
        ws_basket.cell(row=row, column=5, value="YES" if metrics.get('is_worst_performer', False) else "NO")
        
        # Get volatility from risk metrics
        risk_data = analysis.risk_metrics.get(symbol, {})
        volatility = risk_data.get('volatility_annualized', 0)
        ws_basket.cell(row=row, column=6, value=f"{volatility:.2f}%")
        row += 1
    
    wb.save(file_path)

def create_powerpoint_report(analysis: FCNAnalysisResult, file_path: Path):
    """Create PowerPoint presentation for FCN basket analysis"""
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "FCN Basket Investment Analysis"
    subtitle.text = f"Analysis Date: {analysis.created_at.strftime('%B %d, %Y')}\nBasket: {', '.join([s.symbol for s in analysis.stocks_info])}"
    
    # Executive Summary slide
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Executive Summary"
    
    params = analysis.request_params.fcn_params
    basket_metrics = analysis.risk_metrics.get('basket_metrics', {})
    
    summary_text = f"""
    FCN Basket Structure:
    • Basket Type: {params.basket_type.replace('_', ' ').title()}
    • Coupon Rate: {params.coupon_rate}% per annum
    • Face Value: ${params.face_value:,.2f}
    • Maturity: {params.maturity_months} months
    • Barrier Levels: {params.knock_out_barrier_pct}% KO / {params.knock_in_barrier_pct}% KI
    
    Key Findings:
    • {len(analysis.stocks_info)} stocks in basket
    • Expected Payoff: ${basket_metrics.get('expected_payoff', 0):,.2f}
    • Knock-Out Probability: {basket_metrics.get('knock_out_probability', 0):.2f}%
    • Knock-In Probability: {basket_metrics.get('knock_in_probability', 0):.2f}%
    • Most Frequent Worst Performer: {basket_metrics.get('most_frequent_worst', 'N/A')}
    """
    
    content.text = summary_text
    
    # Basket Analysis slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Basket Composition & Performance"
    
    stock_text = "Current Basket Status:\n\n"
    for stock in analysis.stocks_info:
        symbol = stock.symbol
        ref_price = params.reference_prices.get(symbol, stock.current_price)
        performance = ((stock.current_price / ref_price) - 1) * 100
        
        stock_text += f"• {symbol} ({stock.exchange})\n"
        stock_text += f"  Current: ${stock.current_price:.2f} | Reference: ${ref_price:.2f}\n"
        stock_text += f"  Performance vs Reference: {performance:+.2f}%\n"
        
        # Check if worst performer
        metrics = analysis.fcn_metrics.get(symbol, {})
        if metrics.get('is_worst_performer', False):
            stock_text += f"  Status: WORST PERFORMER\n"
        else:
            stock_text += f"  Status: Outperforming\n"
        stock_text += "\n"
    
    content.text = stock_text
    
    # Risk Analysis slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Risk Assessment"
    
    risk_text = "Basket Risk Profile:\n\n"
    
    # Basket level risks
    risk_text += f"Basket Metrics:\n"
    risk_text += f"• Expected Payoff: ${basket_metrics.get('expected_payoff', 0):,.2f}\n"
    risk_text += f"• Average Redemption: {basket_metrics.get('avg_redemption_month', 0):.1f} months\n"
    risk_text += f"• VaR 95%: ${basket_metrics.get('var_95', 0):,.2f}\n"
    risk_text += f"• VaR 99%: ${basket_metrics.get('var_99', 0):,.2f}\n\n"
    
    # Individual stock risks
    risk_text += "Individual Stock Risk:\n"
    for symbol in [s.symbol for s in analysis.stocks_info]:
        risk_data = analysis.risk_metrics.get(symbol, {})
        risk_text += f"• {symbol}:\n"
        risk_text += f"  Volatility: {risk_data.get('volatility_annualized', 0):.2f}%\n"
        risk_text += f"  Max Drawdown: {risk_data.get('max_drawdown', 0):.2f}%\n"
        risk_text += f"  Sharpe Ratio: {risk_data.get('sharpe_ratio', 0):.2f}\n\n"
    
    content.text = risk_text
    
    # Scenario Analysis slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Scenario Analysis"
    
    scenario_text = "Investment Scenarios:\n\n"
    for scenario_name, scenario_data in analysis.scenario_analysis.items():
        scenario_text += f"{scenario_name.replace('_', ' ').title()}:\n"
        scenario_text += f"• Basket Performance: {scenario_data.get('basket_performance', 0):+.2f}%\n"
        scenario_text += f"• Total Payoff: ${scenario_data.get('payoff', 0):,.2f}\n"
        scenario_text += f"• Total Return: {scenario_data.get('total_return', 0):+.2f}%\n"
        scenario_text += f"• Worst Performer: {scenario_data.get('worst_performer', 'N/A')}\n"
        scenario_text += f"• Redemption Type: {scenario_data.get('redemption_type', 'N/A').replace('_', ' ')}\n\n"
    
    content.text = scenario_text
    
    # Investment Recommendation slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Investment Recommendation"
    
    # Generate basic recommendation based on metrics
    ko_prob = basket_metrics.get('knock_out_probability', 0)
    ki_prob = basket_metrics.get('knock_in_probability', 0)
    expected_return = (basket_metrics.get('expected_payoff', params.face_value) / params.face_value - 1) * 100
    
    if ko_prob > 30 and ki_prob < 20:
        recommendation = "ATTRACTIVE"
        reason = "High early redemption probability with low downside risk"
    elif ki_prob > 40:
        recommendation = "HIGH RISK"
        reason = "Elevated knock-in probability increases capital risk"
    else:
        recommendation = "MODERATE"
        reason = "Balanced risk-return profile for coupon enhancement"
    
    recommendation_text = f"""
    Investment Assessment: {recommendation}
    
    Key Factors:
    • Expected Return: {expected_return:+.2f}%
    • Early Redemption Probability: {ko_prob:.2f}%
    • Capital Risk Probability: {ki_prob:.2f}%
    • Coupon Enhancement: {params.coupon_rate}% p.a.
    
    Rationale:
    {reason}
    
    Risk Considerations:
    • Performance linked to worst-performing stock
    • Barrier monitoring may trigger early events
    • Monthly coupon provides income enhancement
    • Capital protection contingent on barrier levels
    """
    
    content.text = recommendation_text
    
    prs.save(file_path)

@api_router.post("/generate-report")
async def generate_report(request: ReportRequest, background_tasks: BackgroundTasks):
    """Generate Excel or PowerPoint report"""
    try:
        # Get analysis data
        analysis_data = await db.fcn_analyses.find_one({"id": request.analysis_id})
        if not analysis_data:
            raise HTTPException(status_code=404, detail="Analysis not found")
        
        analysis = FCNAnalysisResult(**analysis_data)
        
        # Generate filename
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        symbols_str = "_".join([s.symbol for s in analysis.stocks_info])
        
        if request.report_type == "excel":
            filename = f"FCN_Analysis_{symbols_str}_{timestamp}.xlsx"
            file_path = REPORTS_DIR / filename
            create_excel_report(analysis, file_path)
        elif request.report_type == "powerpoint":
            filename = f"FCN_Presentation_{symbols_str}_{timestamp}.pptx"
            file_path = REPORTS_DIR / filename
            create_powerpoint_report(analysis, file_path)
        else:
            raise HTTPException(status_code=400, detail="Invalid report type")
        
        return {"filename": filename, "download_url": f"/api/download/{filename}"}
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Report generation error: {str(e)}")

@api_router.get("/download/{filename}")
async def download_report(filename: str):
    """Download generated report"""
    file_path = REPORTS_DIR / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    
    return FileResponse(
        path=file_path,
        filename=filename,
        media_type='application/octet-stream'
    )

# Include the router in the main app
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
